import os
import json
import requests
from datetime import datetime
import google.genai as genai
from gemini_config import get_gemini_settings
from pptx import Presentation
from pptx.util import Inches, Pt

from speak import speak
from voice_input import VoiceListener

class PresentationAssistant:
    def __init__(self):
        self.api_key, self.model = get_gemini_settings()
        if not self.api_key:
            print("Warning: GOOGLE_GEMINI_KEY not set. Presentation features will be limited.")
            self.client = None
        else:
            self.client = genai.Client(api_key=self.api_key)
        self.voice = VoiceListener(calibrate=False)
        self.presentations_folder = "presentations"
        self.current_presentation = None

        if not os.path.exists(self.presentations_folder):
            os.makedirs(self.presentations_folder)

    def listen(self):
        if not self.voice.available:
            print("Microphone support is unavailable. Please type your command instead.")
            return input("Your command: ").strip().lower()
        return self.voice.listen()

    def generate_content(self, title):
        if not self.client:
            return f"Presentation outline for {title}: Introduction, Key Points, Analysis, Conclusion, Summary."
        try:
            chat = self.client.chats.create(model=self.model)
            response = chat.send_message(
                f"Create an outline for a presentation titled '{title}'. Include 5 main points with brief descriptions.",
                config={"max_output_tokens": 500, "temperature": 0.7},
            )
            result = getattr(response, "text", "") or ""
            return result.strip() or f"Presentation outline for {title}: Introduction, Key Points, Analysis, Conclusion, Summary."
        except Exception as e:
            error_message = f"Unexpected error: {str(e)}"
            print(error_message)
            return error_message
        

    def search_image(self, query):
        try:
            search_url = "https://www.google.com/search?tbm=isch&q=" + requests.utils.quote(query)
            response = requests.get(search_url, timeout=10)
            response.raise_for_status()
            return response.url
        except Exception as exc:
            print(f"Image search fallback failed: {exc}")
            return None

    def create_presentation(self):
        speak("What should be the title of the presentation?")
        title = self.listen()
        if not title:
            return "I couldn't understand the title. Please try again."

        speak(f"Creating a presentation on {title}. This may take a few moments.")
        
        # Generate content
        content = self.generate_content(title)
        slides = self.parse_content(content)
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title
        subtitle_shape.text = f"Created by AVA on {datetime.now().strftime('%Y-%m-%d')}"
        
        # Content slides
        for slide_content in slides:
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            
            title_shape = slide.shapes.title
            title_shape.text = slide_content['title']
            
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = "\n".join(slide_content['points'])
            
            # Add image
            img_url = self.search_image(slide_content['title'])
            if img_url:
                img_data = requests.get(img_url).content
                img_path = f"temp_image_{slide_content['title'].replace(' ', '_')}.jpg"
                with open(img_path, 'wb') as handler:
                    handler.write(img_data)
                slide.shapes.add_picture(img_path, Inches(5), Inches(2), height=Inches(3))
                os.remove(img_path)
        
        # Save presentation
        filename = f"{self.presentations_folder}/{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        prs.save(filename)
        
        self.current_presentation = filename
        return f"Created presentation '{title}' with {len(slides)} slides. You can find it at {filename}"

    def parse_content(self, content):
        slides = []
        current_slide = None
        for line in content.split('\n'):
            if line.startswith("Slide"):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {"title": "", "points": []}
            elif line.strip().startswith("Title:"):
                current_slide["title"] = line.split("Title:")[1].strip()
            elif line.strip().startswith("- "):
                current_slide["points"].append(line.strip()[2:])
        if current_slide:
            slides.append(current_slide)
        return slides

    def process_command(self, command):
        try:
            if "presentation" in command.lower():
                title = command.lower().replace("presentation", "").strip()
                if not title:
                    return "Please provide a presentation title."
                content = self.generate_content(title)
                if content:
                    return content
                return "Unable to generate presentation content."
            return "Please specify a presentation title."
        except Exception as e:
            return f"Error processing command: {str(e)}"

def main():
    try:
        assistant = PresentationAssistant()
        print("Presentation Assistant initialized successfully!")
        speak("Hello! I'm your AI presentation assistant. How can I help you today?")

        while True:
            print("\nWaiting for your command... (say 'help' for available commands)")
            command = assistant.listen()
            
            if command:
                if any(word in command for word in ['quit', 'exit', 'bye', 'goodbye']):
                    speak("Goodbye! Have a great day!")
                    break
                    
                response = assistant.process_command(command)
                print("\nResponse:", response)
                speak(response)
                
            else:
                print("No command detected. Please try again.")

    except KeyboardInterrupt:
        print("\nProgram terminated by user.")
    except Exception as e:
        print(f"An error occurred: {str(e)}")
    finally:
        print("\nThank you for using the AI Presentation Assistant!")

if __name__ == "__main__":
    main()
